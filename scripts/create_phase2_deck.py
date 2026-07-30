"""Create the NeuroX Phase 2 GhostBusters presentation deck.

The deck is intentionally built from project facts and the supplied Phase 2
guidelines. It adds a fade transition to every slide using PresentationML so
the output remains a native, editable PowerPoint file.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as Shape
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "GhostBusters_NeuroX_Phase_2_Presentation.pptx"
NOTES = ROOT / "GhostBusters_NeuroX_Phase_2_Presenter_Notes.md"
LOGO = ROOT / "static" / "assets" / "Ghostopslogo.png"

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(8, 18, 31)
PANEL = RGBColor(16, 38, 57)
PANEL_2 = RGBColor(20, 49, 70)
BLUE = RGBColor(65, 190, 232)
CYAN = RGBColor(72, 222, 216)
PURPLE = RGBColor(167, 109, 241)
GOLD = RGBColor(246, 191, 66)
GREEN = RGBColor(91, 204, 126)
RED = RGBColor(242, 105, 105)
WHITE = RGBColor(244, 248, 252)
MUTED = RGBColor(174, 196, 210)


def add_transition(slide) -> None:
    """Add a medium fade transition; works in desktop PowerPoint."""
    transition = parse_xml(
        '<p:transition {} spd="med" advClick="1"><p:fade/></p:transition>'.format(nsdecls("p"))
    )
    slide.element.append(transition)


def set_background(slide) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    # Architectural accent lines, deliberately subtle.
    for y, color, width in [(0.17, BLUE, 0.018), (7.26, CYAN, 0.012)]:
        line = slide.shapes.add_shape(Shape.RECTANGLE, 0, Inches(y), W, Inches(width))
        line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()
    glow = slide.shapes.add_shape(Shape.OVAL, Inches(10.9), Inches(-0.8), Inches(3.4), Inches(3.4))
    glow.fill.solid(); glow.fill.fore_color.rgb = PANEL_2; glow.fill.transparency = 28; glow.line.fill.background()


def text(slide, value, x, y, w, h, *, size=18, color=WHITE, bold=False,
         font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = value
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def title(slide, heading, subheading=None, number=None):
    if number is not None:
        text(slide, f"{number:02d}", 0.56, 0.36, 0.52, 0.32, size=10, color=CYAN, bold=True)
    text(slide, heading, 0.56, 0.62, 11.7, 0.6, size=27, color=WHITE, bold=True)
    if subheading:
        text(slide, subheading, 0.58, 1.22, 11.4, 0.36, size=11.5, color=MUTED)


def card(slide, x, y, w, h, heading, body, accent=BLUE, number=None):
    shape = slide.shapes.add_shape(Shape.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = PANEL; shape.line.color.rgb = accent; shape.line.width = Pt(1.2)
    band = slide.shapes.add_shape(Shape.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    band.fill.solid(); band.fill.fore_color.rgb = accent; band.line.fill.background()
    if number is not None:
        badge = slide.shapes.add_shape(Shape.OVAL, Inches(x + 0.23), Inches(y + 0.22), Inches(0.36), Inches(0.36))
        badge.fill.solid(); badge.fill.fore_color.rgb = accent; badge.line.fill.background()
        text(slide, str(number), x + 0.23, y + 0.255, 0.36, 0.15, size=8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        tx = x + 0.72
    else:
        tx = x + 0.27
    text(slide, heading, tx, y + 0.22, w - (tx - x) - 0.25, 0.32, size=13, color=accent, bold=True)
    text(slide, body, x + 0.27, y + 0.68, w - 0.52, h - 0.82, size=10.5, color=WHITE)


def pill(slide, value, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(Shape.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.33))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.fill.transparency = 12; shape.line.color.rgb = color
    text(slide, value, x + 0.08, y + 0.075, w - 0.16, 0.15, size=8.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def arrow(slide, x, y, w=0.45, color=BLUE):
    shp = slide.shapes.add_shape(Shape.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.32))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()


def footer(slide, page, label="GhostBusters · NeuroX Phase 2"):
    text(slide, label, 0.58, 7.06, 4.2, 0.18, size=8, color=MUTED)
    text(slide, str(page), 12.15, 7.04, 0.5, 0.2, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def slide(prs, heading, subheading=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s); title(s, heading, subheading, len(prs.slides))
    add_transition(s)
    return s


def build() -> Presentation:
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H

    # 1. Cover
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_background(s); add_transition(s)
    if LOGO.exists(): s.shapes.add_picture(str(LOGO), Inches(0.64), Inches(0.58), width=Inches(2.0))
    pill(s, "NEUROX 1.0 · PHASE 2", 0.7, 1.68, 2.3, GOLD)
    text(s, "GhostBusters", 0.66, 2.13, 8.2, 0.75, size=39, color=WHITE, bold=True)
    text(s, "A controlled autonomous FinOps agent for safe cloud-cost remediation", 0.68, 2.98, 8.7, 0.48, size=19, color=CYAN)
    text(s, "From a high-level objective to evidence, policy, human review, and an auditable remediation proposal.", 0.7, 3.63, 8.8, 0.42, size=13, color=MUTED)
    for x, label, color in [(0.7, "Goal-directed autonomy", PURPLE), (3.15, "Multi-tool evidence", BLUE), (5.48, "Human-controlled action", GREEN)]:
        pill(s, label, x, 4.42, 2.05, color)
    text(s, "Team: [Add team name / members]", 0.7, 6.55, 5.0, 0.25, size=10, color=MUTED)
    footer(s, 1)

    # 2. Guideline alignment
    s = slide(prs, "Built to answer the Phase 2 rubric", "Every core requirement is visible in the product and deliberately highlighted in this presentation.")
    items = [
        ("25%", "B2B impact & viability", "FinOps cost-control workflow for cloud and Terraform teams", GOLD),
        ("25%", "Autonomous reasoning", "Goal → plan → selective tools → branch/retry → decision", PURPLE),
        ("20%", "Technical architecture", "FastAPI, Pydantic contracts, adapters, PostgreSQL, Redis, policy", BLUE),
        ("15%", "Human in the loop", "Approve, reject, modify, request evidence, add context, waive", GREEN),
        ("15%", "Live demo", "A controlled, explainable end-to-end agent run", RGBColor(244, 137, 68)),
    ]
    for i, (pct, h, b, c) in enumerate(items):
        x = 0.7 + (i % 3) * 4.15; y = 1.83 + (i // 3) * 2.2
        card(s, x, y, 3.75, 1.72, h, b, c)
        text(s, pct, x + 2.72, y + 0.18, 0.7, 0.3, size=18, color=c, bold=True, align=PP_ALIGN.RIGHT)
    card(s, 0.7, 6.1, 12.0, 0.6, "Presentation principle", "We demonstrate genuine, bounded autonomy—not a chatbot, a fixed pipeline, a single API call, or a fake demo.", RED)
    footer(s, 2)

    # 3. Problem / B2B
    s = slide(prs, "The B2B problem: cloud waste is easy to see, hard to change safely", "Teams need savings opportunities with enough evidence and governance to act—not another noisy dashboard.")
    card(s, 0.7, 1.9, 3.75, 3.55, "The friction", "Low utilization alone cannot justify a change. Ownership, dependencies, production risk, business context, pricing, and recent engineering activity may tell a different story.", RED, 1)
    card(s, 4.78, 1.9, 3.75, 3.55, "The gap", "FinOps findings often become manual investigations spread across Terraform, cloud consoles, GitHub, Jira, and people. Decisions are slow and hard to audit.", GOLD, 2)
    card(s, 8.86, 1.9, 3.75, 3.55, "GhostBusters", "A controlled agent converts a broad objective into a traceable recommendation: what it investigated, why it chose each tool, what it found, and where a human must decide.", CYAN, 3)
    pill(s, "B2B value: safer savings discovery + lower investigation effort + auditable decisions", 1.15, 6.0, 11.0, GREEN)
    footer(s, 3)

    # 4. autonomy
    s = slide(prs, "What makes GhostBusters an autonomous agent?", "It decides what to investigate next from the goal, resource state, evidence quality, and safety constraints.")
    steps = [
        ("Goal", "Interpret an objective", PURPLE), ("Plan", "Choose questions & tools", BLUE),
        ("Observe", "Collect structured evidence", CYAN), ("Reason", "Compare alternatives", GOLD),
        ("Branch", "Retry / escalate / abstain", RED), ("Deliver", "Proposal + decision trail", GREEN),
    ]
    for i, (h, b, c) in enumerate(steps):
        x = 0.45 + i * 2.1
        card(s, x, 2.25, 1.7, 1.75, h, b, c, i + 1)
        if i < len(steps) - 1: arrow(s, x + 1.74, 2.95, 0.28, c)
    card(s, 0.82, 4.65, 5.65, 1.25, "Dynamic rather than fixed", "Production/destructive signals skip normal optimization. Missing evidence triggers a request for evidence or abstention. Active dependencies block remediation. New human context re-runs the evaluation.", PURPLE)
    card(s, 6.85, 4.65, 5.65, 1.25, "Bounded rather than black-box", "Optional Gemini may propose an explanation or registered read-only tool. Deterministic validation, verifier checks, policy, and human approval retain final authority.", CYAN)
    footer(s, 4)

    # 5 architecture
    s = slide(prs, "Clean architecture: orchestration separated from evidence and authority", "The agent is testable without live infrastructure; integrations are narrow adapters, not hidden logic.")
    layers = [
        ("GhostOps UI", "Overview · Goals · PR Reviews · Cloud Hunt · Approvals · Audit", BLUE),
        ("FastAPI + Pydantic", "Authenticated API, role checks, validation, webhooks, readiness", CYAN),
        ("Workflow services", "WorkflowService · CloudHuntService · outcome verification · assistant", PURPLE),
        ("Decision engine", "Planner · investigator · conflicts · alternatives · verifier · confidence · policy", GOLD),
        ("Adapter + storage boundary", "GitHub · Jira · AWS · Terraform · PostgreSQL · Redis", GREEN),
    ]
    for i, (h, b, c) in enumerate(layers):
        y = 1.7 + i * 0.94
        card(s, 1.1, y, 11.15, 0.7, h, b, c, i + 1)
    text(s, "Why it matters: tool calls, policy, storage, and UI can evolve independently—while the safety contract stays explicit.", 1.15, 6.55, 11.0, 0.25, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 5)

    # 6 tools
    s = slide(prs, "Multi-tool integration: the agent selects evidence, not just APIs", "The guideline asks for 3+ external tools/APIs and intelligent tool choice. GhostBusters implements both.")
    tools = [
        ("Terraform", "Change, environment, destructive/replacement signals", BLUE),
        ("GitHub", "PR files, reviews, commits, CODEOWNERS", CYAN),
        ("Jira", "Business state, owner, decommission/migration signals", PURPLE),
        ("AWS", "Inventory, CloudWatch utilization, pricing provenance", GOLD),
        ("Dependencies", "Active downstream protections", GREEN),
        ("Pricing", "Savings estimate with source/provenance", RGBColor(244, 137, 68)),
    ]
    for i, (h, b, c) in enumerate(tools):
        x = 0.7 + (i % 3) * 4.15; y = 1.78 + (i // 3) * 1.72
        card(s, x, y, 3.75, 1.32, h, b, c, i + 1)
    card(s, 0.7, 5.48, 12.0, 0.82, "Example of agent selection", "A GitHub-related goal activates PR context, ownership, activity and review tools. An AWS-related goal activates inventory, CloudWatch, tags and pricing. The planner records selected and skipped tools with reasons.", CYAN)
    footer(s, 6)

    # 7 workflow
    s = slide(prs, "Autonomous workflow in action: Terraform PR Review", "An end-to-end loop from a high-level goal to a human-reviewed remediation proposal.")
    flow = ["Goal / webhook", "Parse change", "Plan evidence", "Collect + retry", "Detect conflicts", "Policy + confidence", "Human decision"]
    for i, item in enumerate(flow):
        x = 0.35 + i * 1.84
        card(s, x, 2.05, 1.48, 1.08, item, "", [PURPLE, BLUE, CYAN, CYAN, GOLD, RED, GREEN][i], i + 1)
        if i < len(flow) - 1: arrow(s, x + 1.5, 2.45, 0.28, BLUE)
    card(s, 0.72, 4.08, 3.8, 1.35, "Safe result", "Eligible staging right-sizing recommendation; approval produces a simulated PR by default, or a narrowly validated real PR only when explicitly enabled.", GREEN)
    card(s, 4.77, 4.08, 3.8, 1.35, "Ambiguous result", "Conflicting Jira and Git signals lower confidence and route the case to request evidence / human context rather than pretending certainty.", GOLD)
    card(s, 8.82, 4.08, 3.8, 1.35, "Unsafe result", "Production, destructive changes, active dependencies, unknown ownership, or critical missing evidence block remediation.", RED)
    footer(s, 7)

    # 8 failure recovery
    s = slide(prs, "Adaptive behavior and failure recovery—not a scripted pipeline", "The agent changes course when tools fail or evidence changes the safety picture.")
    cases = [
        ("Missing evidence", "Tool timeout/unavailability is recorded as unavailable evidence. Bounded retry/backoff runs. The agent requests evidence or abstains; it never invents data.", RED),
        ("Conflicting evidence", "Recent Git activity can contradict a completed Jira issue. Conflict detection penalizes confidence and routes to human context.", GOLD),
        ("Protection signal", "Production or an active dependency changes the plan: normal optimization stops and policy blocks remediation.", PURPLE),
        ("Human feedback", "A reviewer can add context, modify, request new evidence, reject, waive, revoke approval, or reopen. The case state and audit trail update.", GREEN),
    ]
    for i, (h, b, c) in enumerate(cases):
        x = 0.75 + (i % 2) * 6.0; y = 1.9 + (i // 2) * 2.12
        card(s, x, y, 5.55, 1.65, h, b, c, i + 1)
    pill(s, "Demo proof point: run Safe → Conflicting → Missing Evidence to show three different agent branches", 1.1, 6.36, 11.1, CYAN)
    footer(s, 8)

    # 9 human
    s = slide(prs, "Human-in-the-loop: autonomy with meaningful control", "GhostBusters prepares the decision; authorized people own the decision to proceed.")
    card(s, 0.7, 1.85, 3.75, 3.68, "Before action", "The system surfaces evidence, conflicts, alternatives, confidence, verifier findings, policy result, estimated savings, and explicit risks. Human approval is mandatory for remediation.", BLUE, 1)
    card(s, 4.8, 1.85, 3.75, 3.68, "At the decision", "Reviewers can approve, reject, modify the recommendation, request evidence, add context, create a waiver, revoke approval, or reopen a case. Roles and organization boundaries are enforced.", GREEN, 2)
    card(s, 8.9, 1.85, 3.75, 3.68, "After the decision", "The result is an auditable proposal—not infrastructure mutation. Outcome verification can record deployment confirmation, observed savings, health signals, partial outcomes, or regressions.", GOLD, 3)
    pill(s, "Non-negotiable safety boundary: no terraform apply · no automatic merge · no cloud mutation", 1.0, 6.05, 11.35, RED)
    footer(s, 9)

    # 10 transparency
    s = slide(prs, "Transparent decision trail: judges can inspect the agent’s work", "The platform is designed to answer: Why this action? Why these tools? What evidence changed the outcome?")
    checkpoints = [
        ("Goal received", "Objective + source", PURPLE), ("Plan created", "Questions, selected/skipped tools", BLUE),
        ("Tools executed", "Input, output, freshness, reliability", CYAN), ("Reasoning", "Conflicts, alternatives, confidence", GOLD),
        ("Safety gate", "Verifier + Rego/Conftest policy", RED), ("Human decision", "Actor, action, comment, version", GREEN),
    ]
    for i, (h, b, c) in enumerate(checkpoints):
        x = 0.7 + (i % 3) * 4.15; y = 1.8 + (i // 3) * 1.78
        card(s, x, y, 3.75, 1.38, h, b, c, i + 1)
    text(s, "Where to show it live", 0.78, 5.6, 2.3, 0.25, size=12, color=CYAN, bold=True)
    text(s, "Technical Audit → tool executions / planning decisions / policy results     •     Activity Log → shared correlation trail     •     Review details → evidence and rationale", 0.78, 5.98, 11.5, 0.45, size=11.2, color=WHITE)
    footer(s, 10)

    # 11 demo
    s = slide(prs, "Live demonstration plan: show the reasoning, not only the UI", "A seven-minute sequence aligned with the guideline’s core test.")
    demo = [
        ("0:00–0:45", "Set the objective", "Explain the business goal and the agent’s safety boundary.", PURPLE),
        ("0:45–2:30", "Safe case", "Start Safe; show tool selection, evidence, recommendation, and approval proposal.", GREEN),
        ("2:30–4:20", "Conflict branch", "Start Conflicting; explain why Git vs Jira disagreement changes the outcome.", GOLD),
        ("4:20–5:40", "Failure branch", "Start Missing Evidence; show retry/unavailable evidence and safe escalation or abstention.", RED),
        ("5:40–7:00", "Trace & close", "Open Technical Audit and Activity Log; connect actions to the rubric.", CYAN),
    ]
    for i, (t, h, b, c) in enumerate(demo):
        y = 1.65 + i * 0.95
        text(s, t, 0.75, y + 0.18, 1.0, 0.22, size=10, color=c, bold=True)
        card(s, 1.85, y, 10.75, 0.7, h, b, c, i + 1)
    pill(s, "Use a controlled real GitHub PR only if your credentials/configuration are ready; otherwise label fixture data honestly.", 0.85, 6.52, 11.7, GOLD)
    footer(s, 11)

    # 12 trust / stack
    s = slide(prs, "Technical credibility: allowed stack, production-aware controls", "Built in Python with conventional services and explicit safety engineering—not an agent-as-a-service or no-code wrapper.")
    card(s, 0.7, 1.78, 3.75, 3.75, "Application stack", "Python · FastAPI · Pydantic · JavaScript frontend · PostgreSQL · Redis · Docker Compose · Render deployment blueprint", BLUE, 1)
    card(s, 4.8, 1.78, 3.75, 3.75, "AI + policy", "Optional Gemini / mock provider for constrained planning and explanation; deterministic fallback; Rego policy via Conftest with safe Python fallback", PURPLE, 2)
    card(s, 8.9, 1.78, 3.75, 3.75, "Operational controls", "Auth/RBAC · sessions/CSRF · secret redaction · idempotency · row versions · retries/timeouts · allowlists · readiness checks · migrations", GREEN, 3)
    text(s, "The architecture makes the agent inspectable: integration calls are recorded, decisions are explainable, and external mutation remains deliberately constrained.", 0.9, 6.05, 11.6, 0.38, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 12)

    # 13 rubric summary
    s = slide(prs, "What the judges should remember", "GhostBusters is an autonomous investigation agent with bounded action and full accountability.")
    summary = [
        ("Autonomy", "It turns broad objectives into a dynamic investigation plan and branches on what it learns.", PURPLE),
        ("Tools", "It selects and interprets multiple external evidence sources rather than calling one API.", CYAN),
        ("Control", "It elevates ambiguity and gives humans meaningful, role-governed choices.", GREEN),
        ("Trust", "It exposes evidence, conflicts, policy, confidence, retries, and decisions in an audit trail.", GOLD),
        ("Impact", "It shortens the path from a cost signal to a safe, reviewable remediation proposal.", BLUE),
    ]
    for i, (h, b, c) in enumerate(summary):
        y = 1.55 + i * 0.88
        card(s, 1.0, y, 11.35, 0.65, h, b, c, i + 1)
    footer(s, 13)

    # 14 close
    s = slide(prs, "GhostBusters", "Autonomous investigation. Human-controlled remediation. Evidence you can trust.")
    if LOGO.exists(): s.shapes.add_picture(str(LOGO), Inches(5.45), Inches(1.55), width=Inches(2.4))
    text(s, "Questions?", 3.8, 3.45, 5.7, 0.65, size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    pill(s, "Goal-directed · multi-tool · adaptive · auditable · human-controlled", 2.0, 4.48, 9.3, CYAN)
    text(s, "Live demo: GhostOps dashboard → PR Review / Cloud Hunt → Technical Audit", 2.05, 5.3, 9.2, 0.35, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 14)
    return prs


def write_notes() -> None:
    NOTES.write_text("""# GhostBusters — NeuroX Phase 2 Presenter Notes

## Suggested opening
“GhostBusters is a controlled autonomous FinOps agent. It does not blindly optimize cloud resources. It turns a goal into an evidence-led investigation, branches when conditions change, and gives people the final authority to act.”

## Demonstration rule
Be accurate about data source mode. If a live GitHub/AWS/Jira integration is not configured, say that the fixture is a controlled demonstration input. Do not claim a cloud mutation, Terraform apply, or real pricing when it did not happen.

## High-value talking points
- Slide 4: Explain the actual branch conditions: production, destructive action, missing evidence, active dependencies, conflicts, and human context.
- Slide 6: Point out that the planner records why it selected or skipped each registered tool.
- Slide 8: This is the strongest answer to “why is this not a fixed pipeline?” Run the conflicting and missing-evidence scenarios live.
- Slide 9: Autonomy is bounded by verifier, policy, RBAC, and a required human decision.
- Slide 10: Open Technical Audit and Activity Log. Show a real correlation trail rather than just describing it.

## Likely judge questions
1. **What makes this autonomous?** The goal-driven planner chooses a relevant subset of tools, changes course based on evidence and safety signals, retries bounded external calls, and can request information or abstain.
2. **How does it choose tools?** The plan considers the goal, Terraform change, environment, and registered capabilities. For example, GitHub terms select PR context/ownership/reviews; AWS terms select inventory, CloudWatch, tags, and pricing.
3. **How does it recover?** It records unavailable evidence, uses bounded retry/backoff, lowers confidence, and safely requests evidence or abstains instead of fabricating results.
4. **Where is the human in the loop?** A reviewer can approve, reject, modify, request evidence, add context, waive, revoke, and reopen. Remediation always requires approval.
5. **Can it change infrastructure?** No Terraform apply or cloud mutation occurs. Real GitHub PR creation is opt-in, guarded, and still goes through normal engineering/CI/CD deployment.
""", encoding="utf-8")


if __name__ == "__main__":
    deck = build()
    deck.save(OUT)
    write_notes()
    print(OUT)
    print(NOTES)
